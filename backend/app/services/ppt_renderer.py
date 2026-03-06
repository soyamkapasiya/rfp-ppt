"""
ppt_renderer.py – Attractive, professional RFP PowerPoint generator.

Generates a fully-branded, Genspark-AI-Slides-style deck with:
  • Custom color theme (deep navy + gold accent)
  • Title slide with gradient background
  • Section divider slides
  • Content slides with icon-style bullet markers
  • Consistent header bar on every slide
  • Footer with slide number
"""

from __future__ import annotations

import textwrap
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Brand palette ──────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0D, 0x1B, 0x3E)   # deep navy background
NAVY_MID    = RGBColor(0x15, 0x2A, 0x5C)   # slightly lighter navy
GOLD        = RGBColor(0xF5, 0xA6, 0x23)   # warm gold accent
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY  = RGBColor(0xE8, 0xEC, 0xF5)
TEAL        = RGBColor(0x00, 0xB4, 0xD8)   # teal for bullet markers
SLIDE_W     = Inches(13.33)
SLIDE_H     = Inches(7.5)

# Slide-type categories for divider detection
DIVIDER_KEYWORDS = {
    "technical architecture",
    "delivery methodology",
    "team and governance",
    "appendix",
}


def _rgb(r: int, g: int, b: int) -> RGBColor:
    return RGBColor(r, g, b)


def _add_rect(slide, left, top, width, height, fill_color: RGBColor | None = None, line_color: RGBColor | None = None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height,
    )
    shape.line.fill.background()  # no line by default
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text: str, font_size: int,
                 bold=False, color: RGBColor = WHITE, align=PP_ALIGN.LEFT,
                 italic=False, wrap=True) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Times New Roman"


def _set_slide_background(slide, color: RGBColor) -> None:
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_header_bar(slide, title_text: str) -> None:
    """Top navy bar with slide title."""
    _add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), fill_color=NAVY)
    # Gold accent line under header
    _add_rect(slide, 0, Inches(1.1), SLIDE_W, Pt(3), fill_color=GOLD)
    _add_textbox(
        slide,
        Inches(0.35), Inches(0.18),
        Inches(11.5), Inches(0.75),
        title_text,
        font_size=24, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
    )


def _add_footer(slide, slide_num: int, total: int) -> None:
    """Bottom bar with slide number and brand name."""
    _add_rect(slide, 0, Inches(7.05), SLIDE_W, Inches(0.45), fill_color=NAVY)
    _add_textbox(
        slide, Inches(0.3), Inches(7.07),
        Inches(8), Inches(0.35),
        "shivanski technology llp  •  sales@shiavnski.in",
        font_size=9, color=LIGHT_GREY, align=PP_ALIGN.LEFT,
    )
    _add_textbox(
        slide, Inches(11.5), Inches(7.07),
        Inches(1.5), Inches(0.35),
        f"{slide_num} / {total}",
        font_size=9, color=GOLD, align=PP_ALIGN.RIGHT,
    )


def _build_title_slide(prs: Presentation, project_name: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
    _set_slide_background(slide, NAVY)

    # Full-width gold bar at top
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.12), fill_color=GOLD)

    # Navy overlay panel (slight lighter box for title)
    _add_rect(slide, Inches(0.6), Inches(1.4), Inches(12.1), Inches(3.5),
              fill_color=NAVY_MID)

    # "RFP PROPOSAL" label
    _add_textbox(
        slide, Inches(0.9), Inches(1.6), Inches(11), Inches(0.5),
        "REQUEST FOR PROPOSAL  •  PROPOSAL RESPONSE",
        font_size=11, bold=False, color=GOLD, align=PP_ALIGN.LEFT,
        italic=False,
    )

    # Project name
    wrapped = "\n".join(textwrap.wrap(project_name.upper(), width=40))
    _add_textbox(
        slide, Inches(0.9), Inches(2.05), Inches(11), Inches(1.6),
        wrapped,
        font_size=34, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
    )

    # Subtitle
    _add_textbox(
        slide, Inches(0.9), Inches(3.75), Inches(10), Inches(0.7),
        subtitle,
        font_size=15, bold=False, color=TEAL, align=PP_ALIGN.LEFT,
    )

    # Decorative teal bar on left edge
    _add_rect(slide, 0, Inches(1.4), Inches(0.08), Inches(3.5), fill_color=TEAL)

    # Bottom section
    _add_rect(slide, 0, Inches(6.3), SLIDE_W, Inches(1.2), fill_color=NAVY_MID)
    _add_textbox(
        slide, Inches(0.9), Inches(6.45), Inches(11), Inches(0.5),
        "Prepared by shivanski technology llp  •  https://www.shiavnski.in/",
        font_size=10, bold=False, color=GOLD, align=PP_ALIGN.LEFT,
    )
    _add_rect(slide, 0, Inches(7.38), SLIDE_W, Inches(0.12), fill_color=GOLD)


def _build_divider_slide(prs: Presentation, section_title: str, slide_num: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, NAVY)
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill_color=GOLD)
    _add_rect(slide, 0, Inches(7.42), SLIDE_W, Inches(0.08), fill_color=GOLD)

    # Large section number
    _add_textbox(
        slide, Inches(0.7), Inches(2.4), Inches(12), Inches(0.6),
        f"SECTION  {slide_num:02d}",
        font_size=13, bold=True, color=GOLD, align=PP_ALIGN.LEFT,
    )
    # Section title
    _add_textbox(
        slide, Inches(0.7), Inches(2.95), Inches(12), Inches(1.6),
        section_title.upper(),
        font_size=38, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
    )
    _add_rect(slide, Inches(0.7), Inches(4.65), Inches(3.5), Pt(3), fill_color=TEAL)
    _add_footer(slide, slide_num, total)


def _build_content_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    references: list[str],
    slide_num: int,
    total: int,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_background(slide, WHITE)

    # White main area, navy header
    _add_header_bar(slide, title)

    # Body content area background
    _add_rect(slide, 0, Inches(1.13), SLIDE_W, Inches(5.92), fill_color=WHITE)

    # Left accent bar
    _add_rect(slide, 0, Inches(1.13), Inches(0.07), Inches(5.92), fill_color=TEAL)

    # Bullet items
    bullet_y_start = Inches(1.35)
    bullet_spacing = Inches(0.72)
    max_bullets = min(len(bullets), 7)

    for i, bullet in enumerate(bullets[:max_bullets]):
        y_pos = bullet_y_start + i * bullet_spacing
        if y_pos + Inches(0.6) > Inches(6.8):
            break

        # Gold dot / bullet marker
        _add_rect(slide, Inches(0.22), y_pos + Pt(8), Inches(0.16), Inches(0.16),
                  fill_color=GOLD)

        # Bullet text – wrap long bullets
        short = bullet.strip()
        if len(short) > 180:
            short = short[:177] + "…"

        _add_textbox(
            slide,
            Inches(0.5), y_pos,
            Inches(12.5), Inches(0.65),
            short,
            font_size=14, bold=False, color=NAVY, align=PP_ALIGN.LEFT,
        )

    # References section (if any)
    if references:
        _add_rect(slide, Inches(0.07), Inches(6.55), SLIDE_W - Inches(0.07), Pt(1), fill_color=LIGHT_GREY)
        ref_text = "References: " + " | ".join(r[:60] for r in references[:3])
        _add_textbox(
            slide, Inches(0.3), Inches(6.6), Inches(12.5), Inches(0.35),
            ref_text,
            font_size=8, bold=False, color=_rgb(0x88, 0x88, 0x99), align=PP_ALIGN.LEFT,
        )

    _add_footer(slide, slide_num, total)


def render_ppt(slides: list[dict], output_path: str, template_path: str | None = None) -> str:
    """
    Render a professional, branded PPTX from a list of slide dicts.

    Each dict must have: title, bullets (list[str]), references (list[str]).
    An optional 'project_name' key on the first slide sets the title slide headline.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Ensure we have a blank layout (layout index 6 in default pptx)
    # Add blank layout if not enough layouts exist
    while len(prs.slide_layouts) < 7:
        prs.slide_layouts._sldLayoutLst.append(
            prs.slide_layouts[0]._element
        )

    total_slides = len(slides) + 1  # +1 for title slide

    # ── Title slide ───────────────────────────────────────────────────────────
    project_name = "RFP Proposal Response"
    if slides:
        project_name = slides[0].get("project_name") or project_name
    subtitle = "A comprehensive, AI-assisted proposal prepared in response to your Request for Proposal"
    _build_title_slide(prs, project_name, subtitle)

    # ── Content slides ────────────────────────────────────────────────────────
    for idx, spec in enumerate(slides, start=2):
        title  = spec.get("title", f"Slide {idx}")
        bullets = spec.get("bullets") or ["Content to be finalized."]
        references = spec.get("references") or []

        is_divider = any(kw in title.lower() for kw in DIVIDER_KEYWORDS)

        if is_divider:
            _build_divider_slide(prs, title, idx, total_slides)
        else:
            _build_content_slide(prs, title, bullets, references, idx, total_slides)

    out = Path(output_path)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return str(out)
